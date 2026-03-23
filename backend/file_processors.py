"""
File processing utilities for extracting text from various formats
"""
import base64
import io
import logging
from typing import List, Dict, Any
from datetime import datetime

logger = logging.getLogger(__name__)

# Try to import file processing libraries
try:
    import PyPDF2
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False
    logger.warning("PyPDF2 not available - PDF processing disabled")

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("python-docx not available - DOCX processing disabled")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available - PPT processing disabled")

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False
    logger.warning("pandas not available - CSV processing disabled")


def extract_text_from_pdf(file_content: bytes) -> str:
    """Extract text from PDF file"""
    if not PYPDF2_AVAILABLE:
        raise Exception("PDF processing not available - PyPDF2 not installed")

    try:
        pdf_file = io.BytesIO(file_content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {e}")
        raise Exception(f"Failed to extract text from PDF: {str(e)}")


def extract_text_from_docx(file_content: bytes) -> str:
    """Extract text from DOCX file"""
    if not DOCX_AVAILABLE:
        raise Exception("DOCX processing not available - python-docx not installed")

    try:
        docx_file = io.BytesIO(file_content)
        doc = Document(docx_file)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + " | "
                text += "\n"
        return text.strip()
    except Exception as e:
        logger.error(f"Error extracting text from DOCX: {e}")
        raise Exception(f"Failed to extract text from DOCX: {str(e)}")


def extract_text_from_csv(file_content: bytes) -> str:
    """Extract text from CSV file"""
    if not PANDAS_AVAILABLE:
        raise Exception("CSV processing not available - pandas not installed")

    try:
        csv_file = io.BytesIO(file_content)
        df = pd.read_csv(csv_file)
        text = "CSV Data Summary:\n\n"
        text += f"Rows: {len(df)}, Columns: {len(df.columns)}\n"
        text += f"Columns: {', '.join(df.columns.tolist())}\n\n"
        text += "Data Preview (first 10 rows):\n"
        text += df.head(10).to_string()
        return text
    except Exception as e:
        logger.error(f"Error extracting text from CSV: {e}")
        raise Exception(f"Failed to extract text from CSV: {str(e)}")


def extract_text_from_pptx(file_content: bytes) -> str:
    """Extract text from PPTX file"""
    if not PPTX_AVAILABLE:
        raise Exception("PPT processing not available - python-pptx not installed")

    try:
        pptx_file = io.BytesIO(file_content)
        prs = Presentation(pptx_file)
        text = ""
        for i, slide in enumerate(prs.slides):
            text += f"\n--- Slide {i + 1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        logger.error(f"Error extracting text from PPTX: {e}")
        raise Exception(f"Failed to extract text from PPTX: {str(e)}")


def process_file(file_data: Dict[str, Any]) -> Dict[str, Any]:
    """
    Process a file and extract its text content
    file_data: {name, type, content_base64}
    Returns: {name, type, content, text}
    """
    file_name = file_data.get('name', 'unknown')
    file_type = file_data.get('type', '')
    content_b64 = file_data.get('content')

    if not content_b64:
        raise Exception("File content is empty")

    # Decode base64 content
    try:
        file_content = base64.b64decode(content_b64)
    except Exception as e:
        raise Exception(f"Failed to decode file content: {str(e)}")

    # Validate file size (50 MB max)
    if len(file_content) > 50 * 1024 * 1024:
        raise Exception(f"File too large: {len(file_content) / (1024*1024):.2f} MB (max 50 MB)")

    # Extract text based on file type
    text = ""
    if 'pdf' in file_type.lower() or file_name.lower().endswith('.pdf'):
        text = extract_text_from_pdf(file_content)
    elif 'wordprocessingml' in file_type.lower() or file_name.lower().endswith('.docx'):
        text = extract_text_from_docx(file_content)
    elif 'csv' in file_type.lower() or file_name.lower().endswith('.csv'):
        text = extract_text_from_csv(file_content)
    elif 'presentation' in file_type.lower() or file_name.lower().endswith('.pptx'):
        text = extract_text_from_pptx(file_content)
    else:
        raise Exception(f"Unsupported file type: {file_type}. Supported: PDF, DOCX, CSV, PPTX")

    return {
        'name': file_name,
        'type': file_type,
        'content': content_b64,
        'text': text,
        'size': len(file_content)
    }


def process_multiple_files(files: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """
    Process multiple files and extract text from each
    Returns: List of processed file data
    """
    if len(files) > 3:
        raise Exception(f"Too many files: {len(files)} (max 3)")

    results = []
    for file_data in files:
        try:
            processed = process_file(file_data)
            results.append(processed)
        except Exception as e:
            logger.error(f"Failed to process file {file_data.get('name', 'unknown')}: {e}")
            results.append({
                'name': file_data.get('name', 'unknown'),
                'type': file_data.get('type', ''),
                'error': str(e)
            })

    return results
